from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# --- 安全定義區塊 (Safe Fallback) ---
# 為了避免 ImportError，我們手動定義這些常數的數值
# 這樣無論您的 python-pptx 版本為何，程式都能跑
class MSO_ARROWHEAD:
    TRIANGLE = 2  # 標準三角形箭頭

class MSO_CONNECTOR:
    STRAIGHT = 1  # 直線連接符

# 嘗試導入 (如果有的話)，沒有就用上面的預設值
try:
    from pptx.enum.dml import MSO_ARROWHEAD
except ImportError:
    pass

try:
    from pptx.enum.shapes import MSO_CONNECTOR
except ImportError:
    pass
# ------------------------------------

def create_experiment_flowchart():
    prs = Presentation()
    # 設定寬螢幕 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout

    # --- 樣式設定 ---
    c_primary = RGBColor(31, 78, 121)  # 深藍色
    c_accent = RGBColor(192, 80, 77)   # 強調色 (紅/橘)
    c_text_gray = RGBColor(89, 89, 89) # 文字灰
    
    # 輔助函數：建立帶文字的方塊
    def add_box(left, top, width, height, text, color=c_primary, font_size=12, shape_type=MSO_SHAPE.RECTANGLE):
        shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(255, 255, 255)
        
        tf = shape.text_frame
        tf.text = text
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.size = Pt(font_size)
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        tf.paragraphs[0].font.bold = True
        return shape

    def add_text_label(left, top, width, height, text, size=10, color=RGBColor(80, 80, 80)):
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.LEFT
        tf.word_wrap = True

    # --- 1. 標題 ---
    add_box(0.5, 0.3, 12.33, 0.6, "實驗設計流程 (Experimental Workflow)", c_primary, 20)

    # --- 2. 左側：資料源 ---
    add_box(0.5, 1.5, 1.5, 5.0, "資料來源\n(Data Source)", c_text_gray, 14)
    add_box(0.7, 2.5, 1.1, 1.0, "Data 2\n(Training)", RGBColor(127, 127, 127), 11)
    add_box(0.7, 4.5, 1.1, 1.0, "Data 1\n(Testing)", RGBColor(127, 127, 127), 11)

    # --- 3. 中間：四個實驗分支 ---
    
    # 分支 1: Baseline
    add_box(2.5, 1.5, 2.5, 1.0, "1. Baseline 模型\n(Standard HRV)", c_primary, 12, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_label(2.6, 2.5, 2.5, 0.6, "特徵: SDNN, LF, HF...\n基本資料: Age, Sex, BMI")
    
    # 分支 2: External Validation
    add_box(5.5, 1.5, 2.5, 1.0, "2. 外部驗證\n(External Val.)", c_accent, 12, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_label(5.6, 2.5, 2.5, 0.6, "測試集: Data 1\n指標: Specificity, NPV")

    # 分支 3: Scale Only
    add_box(2.5, 3.2, 2.5, 1.0, "3. 比較分析\n(Scale Only)", c_primary, 12, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_label(2.6, 4.2, 2.5, 0.6, "特徵: phq15, bdi...\n(無外部驗證)")

    # 分支 4: Extended HRV
    add_box(2.5, 4.9, 2.5, 1.0, "4. 附加分析\n(Extended HRV)", c_primary, 12, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_label(2.6, 5.9, 2.5, 0.6, "特徵: All HRV + Scale\n(子群組分析)")

    # 分支 5: Full Model
    add_box(5.5, 4.9, 2.5, 1.0, "5. 綜合模型\n(Full Features)", c_primary, 12, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_label(5.6, 5.9, 2.5, 0.6, "訓練: Data1 + Data2\n特徵: 臨床 + 自律神經")

    # --- 4. 右側：評估區 ---
    add_box(8.8, 1.5, 4.0, 5.0, "", RGBColor(230, 230, 230)) # 淺灰背景
    add_box(9.0, 1.7, 3.6, 0.5, "綜合評估 (Evaluation)", RGBColor(31, 78, 121), 14)

    metrics_text = (
        "績效指標 (Metrics):\n"
        "• Accuracy, F1, AUC\n"
        "• Specificity, NPV\n\n"
        "模型解釋性 (XAI):\n"
        "• Random Forest Importance\n"
        "• SHAP / Permutation\n\n"
        "目標 (Goal):\n"
        "• 驗證 HRV 特徵有效性\n"
        "• 比較生理 vs 心理指標"
    )
    
    tb = slide.shapes.add_textbox(Inches(9.0), Inches(2.3), Inches(3.6), Inches(4.0))
    for line in metrics_text.split('\n'):
        p = tb.text_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_after = Pt(6)

    # --- 5. 箭頭繪製函數 (使用安全常數) ---
    def add_arrow(x1, y1, x2, y2):
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, 
            Inches(x1), Inches(y1), Inches(x2), Inches(y2)
        )
        conn.line.width = Pt(2)
        conn.line.color.rgb = RGBColor(150, 150, 150)
        conn.line.end_arrowhead = MSO_ARROWHEAD.TRIANGLE

    # Data 2 -> Models
    add_arrow(1.8, 3.0, 2.5, 2.0) # To Baseline
    add_arrow(1.8, 3.0, 2.5, 3.7) # To Scale
    add_arrow(1.8, 3.0, 2.5, 5.4) # To Extended

    # Data 1 -> Ext Val & Full Model
    add_arrow(1.8, 5.0, 5.5, 2.0) # To Ext Val (Cross)
    add_arrow(1.8, 5.0, 5.5, 5.4) # To Full Model

    # Baseline -> Ext Val
    add_arrow(5.0, 2.0, 5.5, 2.0)

    # Models -> Evaluation
    add_arrow(8.0, 2.0, 8.8, 3.0) # Baseline out
    add_arrow(5.0, 3.7, 8.8, 3.0) # Scale out
    add_arrow(5.0, 5.4, 8.8, 3.0) # Extended out
    add_arrow(8.0, 5.4, 8.8, 3.0) # Full out

    prs.save('experiment_workflow_final.pptx')
    print("PPT generated successfully: experiment_workflow_final.pptx")

if __name__ == "__main__":
    create_experiment_flowchart()